from __future__ import annotations

from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def _add_title(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.5))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(30, 55, 90)


def _add_metric(slide, x, y, title, value):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.35), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(235, 242, 249)
    shape.line.color.rgb = RGBColor(180, 200, 220)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(value); r.font.bold = True; r.font.size = Pt(20); r.font.color.rgb = RGBColor(30, 55, 90)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = title; r2.font.size = Pt(10); r2.font.color.rgb = RGBColor(70, 70, 70)


def export_solutions_to_pptx_bytes(solutions) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "OptiFLUX — Synthèse des simulations")
    total_km = sum(s.metrics.total_km for s in solutions)
    vehicles = sum(len({r.vehicle.id for r in s.routes}) for s in solutions)
    shifts = sum(len(s.routes) for s in solutions)
    occ = sum(s.metrics.driver_useful_occupancy for s in solutions) / len(solutions) if solutions else 0
    _add_metric(slide, 0.7, 1.2, "Jours simulés", len(solutions))
    _add_metric(slide, 3.3, 1.2, "Véhicules cumulés", vehicles)
    _add_metric(slide, 5.9, 1.2, "Postes chauffeurs", shifts)
    _add_metric(slide, 8.5, 1.2, "Km totaux", f"{total_km:.1f}")
    _add_metric(slide, 11.1, 1.2, "Occupation chauffeur", f"{occ:.0%}")
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(12), Inches(3.5))
    tf = tx.text_frame
    tf.text = "Messages clés"
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.size = Pt(18)
    for sol in solutions:
        p = tf.add_paragraph()
        p.text = f"• {sol.day} : {sol.status_message} — {len(sol.routes)} postes, {sol.metrics.total_km:.1f} km, occupation chauffeur {sol.metrics.driver_useful_occupancy:.0%}"
        p.font.size = Pt(12)

    for sol in solutions:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_title(slide, f"Détail — {sol.day}")
        _add_metric(slide, 0.7, 1.1, "Postes", len(sol.routes))
        _add_metric(slide, 3.3, 1.1, "Véhicules", len({r.vehicle.id for r in sol.routes}))
        _add_metric(slide, 5.9, 1.1, "Km", f"{sol.metrics.total_km:.1f}")
        _add_metric(slide, 8.5, 1.1, "Km à vide", f"{sol.metrics.empty_km:.1f}")
        _add_metric(slide, 11.1, 1.1, "Occupation", f"{sol.metrics.driver_useful_occupancy:.0%}")
        tx = slide.shapes.add_textbox(Inches(0.7), Inches(2.7), Inches(12), Inches(3.9))
        tf = tx.text_frame
        tf.text = "Principales tournées"
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.size = Pt(16)
        for r in sol.routes[:12]:
            p = tf.add_paragraph()
            p.text = f"• {r.vehicle.id} / {r.shift.id} : {len(r.unit_ids)} unités, {sum(e.distance_km for e in r.events if e.event_type == 'Trajet'):.1f} km"
            p.font.size = Pt(11)

    out = BytesIO()
    prs.save(out)
    return out.getvalue()


def save_pptx(path: str, solutions) -> None:
    with open(path, "wb") as f:
        f.write(export_solutions_to_pptx_bytes(solutions))
